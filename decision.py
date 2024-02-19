import yfinance as yf
import pandas as pd
import ta
from pptx import Presentation
from pptx.util import Inches

# Step 1: Scrape EUR/INR currency data from Yahoo Finance
ticker_symbol = "EURINR=X"
start_date = "2023-01-01"
end_date = "2024-02-16"
currency_data = yf.download(ticker_symbol, start=start_date, end=end_date)

# Step 2: Perform technical analysis
# Calculate Moving Averages (MA)
currency_data['ma_one_day'] = ta.trend.sma_indicator(currency_data['Close'], window=1)
currency_data['ma_one_week'] = ta.trend.sma_indicator(currency_data['Close'], window=5)

# Calculate Bollinger Bands
bb_indicator = ta.volatility.BollingerBands(currency_data['Close'], window=20)
currency_data['bb_upper'] = bb_indicator.bollinger_hband()
currency_data['bb_lower'] = bb_indicator.bollinger_lband()

# Calculate Commodity Channel Index (CCI)
currency_data['cci_one_day'] = ta.trend.cci(currency_data['High'], currency_data['Low'], currency_data['Close'], window=14)
currency_data['cci_one_week'] = ta.trend.cci(currency_data['High'], currency_data['Low'], currency_data['Close'], window=70)

# Step 3: Make decisions based on technical indicators
decision_one_day = 'BUY' if currency_data['ma_one_day'].iloc[-1] < currency_data['bb_lower'].iloc[-1] and currency_data['cci_one_day'].iloc[-1] > 0 else 'SELL' if currency_data['ma_one_day'].iloc[-1] > currency_data['bb_upper'].iloc[-1] and currency_data['cci_one_day'].iloc[-1] < 0 else 'NEUTRAL'
decision_one_week = 'BUY' if currency_data['ma_one_week'].iloc[-1] < currency_data['bb_lower'].iloc[-1] and currency_data['cci_one_week'].iloc[-1] > 0 else 'SELL' if currency_data['ma_one_week'].iloc[-1] > currency_data['bb_upper'].iloc[-1] and currency_data['cci_one_week'].iloc[-1] < 0 else 'NEUTRAL'

# Print decisions
print("Decision for one day from February 16, 2024:", decision_one_day)
print("Decision for one week from February 16, 2024:", decision_one_week)

# Create a presentation object
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technical Analysis Decisions"

# Add a slide with a table for decision data
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technical Analysis Decision Table"

# Define decision data
decision_data = {
    'Time Frame': ['One Day', 'One Week'],
    'Decision': [decision_one_day, decision_one_week]
}

# Define table dimensions and position
left_inch = Inches(1)
top_inch = Inches(1.5)
width_inch = Inches(3)
height_inch = Inches(1.5)

# Add the table to the slide
table = slide.shapes.add_table(
    rows=len(decision_data) + 1,
    cols=len(decision_data.keys()),
    left=left_inch,
    top=top_inch,
    width=width_inch,
    height=height_inch
).table

# Populate the table with data from the DataFrame
for col_index, column in enumerate(decision_data.keys()):
    cell = table.cell(0, col_index)
    cell.text = column
    for row_index, value in enumerate(decision_data[column]):
        cell = table.cell(row_index + 1, col_index)
        cell.text = str(value)

# Save the presentation
prs.save("technical_analysis_decisions_1.pptx")