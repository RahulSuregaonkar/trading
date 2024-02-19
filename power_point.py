from pptx import Presentation
from pptx.util import Inches
import pandas as pd

# Create a new PowerPoint presentation
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technical Analysis Decision"

# Define the decision data (example data)
decision_data = {
    'Indicator': ['Moving Average', 'Bollinger Bands', 'Commodity Channel Index'],
    'Decision': ['BUY', 'SELL', 'NEUTRAL']
}

# Create a DataFrame from the decision data
df_decision = pd.DataFrame(decision_data)

# Add a slide with the table
slide_layout = prs.slide_layouts[5]  # Use the layout for a title and content slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Individual Indicator Decisions"

# Define table dimensions and position
left_inch = Inches(1)
top_inch = Inches(1.5)
width_inch = Inches(6)
height_inch = Inches(3)

# Add the table to the slide
table = slide.shapes.add_table(
    rows=df_decision.shape[0] + 1,
    cols=df_decision.shape[1],
    left=left_inch,
    top=top_inch,
    width=width_inch,
    height=height_inch
).table

# Populate the table with data from the DataFrame
for col_index, column in enumerate(df_decision.columns):
    cell = table.cell(0, col_index)
    cell.text = column
    for row_index, value in enumerate(df_decision[column]):
        cell = table.cell(row_index + 1, col_index)
        cell.text = str(value)

# Save the presentation
prs.save("technical_analysis_decision.pptx")